#!/usr/bin/env python3
"""
Generate the Eat Jersey Challenge presentation as a PowerPoint file.

Usage:
    python scripts/generate-presentation.py            # 30-min version (~30 slides)
    python scripts/generate-presentation.py --short     # 15-min version (~18 slides)

Output:
    EatJerseyChallenge-Presentation.pptx   (or EatJerseyChallenge-Presentation-Short.pptx)
"""

import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colours ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GREY = RGBColor(0x99, 0x99, 0x99)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
PURPLE = RGBColor(0x80, 0x00, 0x80)

# ── Layout constants ──────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "images", "EJCLogo-Transparent.png")

LOGO_SIZE = Inches(0.55)
LOGO_LEFT = Inches(0.3)
LOGO_TOP = Inches(0.2)

TITLE_LEFT = Inches(1.0)
TITLE_TOP = Inches(0.15)
TITLE_WIDTH = Inches(11.5)
TITLE_HEIGHT = Inches(0.7)

BODY_LEFT = Inches(1.0)
BODY_TOP = Inches(1.1)
BODY_WIDTH = Inches(11.5)
BODY_HEIGHT = Inches(5.8)

FOOTER_TOP = Inches(7.05)
FOOTER_HEIGHT = Inches(0.35)


# ── Helper functions ──────────────────────────────────────────────────────

def add_logo(slide):
    """Add the EJC logo to the top-left of the slide."""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, LOGO_LEFT, LOGO_TOP, height=LOGO_SIZE)


def set_slide_bg(slide, colour):
    """Set the solid background colour of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_textbox(slide, left, top, width, height, text, font_size=18,
                colour=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a simple text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = colour
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    colour=WHITE, font_name="Calibri", bold_first_word=False):
    """Add a bulleted list text box. Items is a list of strings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0

        if bold_first_word and ":" in item:
            before, after = item.split(":", 1)
            run1 = p.add_run()
            run1.text = before + ":"
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = GOLD
            run1.font.bold = True
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = after
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = colour
            run2.font.name = font_name
        else:
            run = p.add_run()
            run.text = f"  •  {item}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = colour
            run.font.name = font_name
    return tf


def add_speaker_notes(slide, notes_text):
    """Set speaker notes for a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_section_header(slide, section_name):
    """Add a small gold section label below the title bar."""
    add_textbox(slide, TITLE_LEFT, Inches(0.85), Inches(4), Inches(0.3),
                section_name, font_size=11, colour=GOLD, bold=True)


def add_placeholder_box(slide, left, top, width, height, label):
    """Add a grey dashed-style placeholder rectangle for screenshots / demos."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)
    # Add label text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {label} ]"
    run.font.size = Pt(14)
    run.font.color.rgb = MEDIUM_GREY
    run.font.italic = True
    run.font.name = "Calibri"
    tf.paragraphs[0].space_before = Pt(4)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_footer(slide, slide_number, total_slides):
    """Add a subtle footer with slide number."""
    add_textbox(slide, Inches(12.0), FOOTER_TOP, Inches(1.0), FOOTER_HEIGHT,
                f"{slide_number} / {total_slides}", font_size=10,
                colour=MEDIUM_GREY, alignment=PP_ALIGN.RIGHT)


def make_title_slide(prs, slide_num, total):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    # Large centered logo
    if os.path.exists(LOGO_PATH):
        logo_big = Inches(2.0)
        slide.shapes.add_picture(
            LOGO_PATH,
            (SLIDE_WIDTH - logo_big) // 2, Inches(0.6),
            height=logo_big)

    add_textbox(slide, Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.2),
                "Building an Interactive Map to\nEat Our Way Through New Jersey",
                font_size=36, colour=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.5),
                "A journey through 564 municipalities — one meal at a time",
                font_size=18, colour=GOLD, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.5),
                "Lessons learned building with GitHub Copilot, Leaflet, and vanilla JavaScript",
                font_size=14, colour=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide, (
        "Welcome everyone. Today I'm going to share the story of how I built "
        "an interactive web map to track our journey to eat at a restaurant in "
        "every single one of New Jersey's 564 municipalities — and what I learned "
        "along the way about AI-assisted development with GitHub Copilot."
    ))
    add_footer(slide, slide_num, total)
    return slide


def make_standard_slide(prs, title, bullets, notes, section=None,
                        slide_num=0, total=0, placeholder_label=None,
                        two_col=None):
    """Create a standard slide with title, optional section label, and bullet content.
    
    two_col: if set, a dict with 'left' and 'right' lists of bullet strings.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)
    add_logo(slide)

    # Title
    add_textbox(slide, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                title, font_size=28, colour=WHITE, bold=True)

    # Gold accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  TITLE_LEFT, Inches(0.8),
                                  Inches(3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    if section:
        add_section_header(slide, section)

    if two_col:
        col_width = Inches(5.3)
        add_bullet_list(slide, BODY_LEFT, BODY_TOP, col_width, BODY_HEIGHT,
                        two_col["left"], bold_first_word=True)
        add_bullet_list(slide, Inches(7.0), BODY_TOP, col_width, BODY_HEIGHT,
                        two_col["right"], bold_first_word=True)
    elif bullets:
        bold_kw = any(":" in b for b in bullets)
        add_bullet_list(slide, BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT,
                        bullets, bold_first_word=bold_kw)

    if placeholder_label:
        ph_top = Inches(4.0) if bullets else BODY_TOP
        ph_height = Inches(3.0) if bullets else Inches(5.5)
        add_placeholder_box(slide, Inches(2.5), ph_top,
                            Inches(8.3), ph_height, placeholder_label)

    if notes:
        add_speaker_notes(slide, notes)

    add_footer(slide, slide_num, total)
    return slide


def make_demo_slide(prs, demo_title, instructions, notes, section,
                    slide_num, total):
    """Create a demo placeholder slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0x12, 0x12, 0x22))
    add_logo(slide)

    add_textbox(slide, TITLE_LEFT, Inches(0.15), TITLE_WIDTH, TITLE_HEIGHT,
                f"DEMO: {demo_title}", font_size=32, colour=GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_placeholder_box(slide, Inches(1.5), Inches(1.5),
                        Inches(10.3), Inches(4.5),
                        "Live Demo")

    add_bullet_list(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(1.0),
                    instructions, font_size=13, colour=LIGHT_GREY)

    add_speaker_notes(slide, notes)
    add_section_header(slide, section)
    add_footer(slide, slide_num, total)
    return slide


def make_thank_you_slide(prs, slide_num, total):
    """Final slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    if os.path.exists(LOGO_PATH):
        logo_big = Inches(1.8)
        slide.shapes.add_picture(
            LOGO_PATH,
            (SLIDE_WIDTH - logo_big) // 2, Inches(0.8),
            height=logo_big)

    add_textbox(slide, Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.8),
                "Thank You!", font_size=40, colour=GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.5),
                "Questions?", font_size=24, colour=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(3.0), Inches(5.2), Inches(7.3), Inches(1.5),
                    ["GitHub: github.com/jrzyshr/eatjerseychallenge",
                     "Live map: [your GitHub Pages URL]"],
                    font_size=16, colour=LIGHT_GREY)

    add_speaker_notes(slide, (
        "Thank the audience. Open up for questions. "
        "Share the GitHub link and live map URL for anyone who wants to explore."
    ))
    add_footer(slide, slide_num, total)
    return slide


# ── Slide definitions ─────────────────────────────────────────────────────

def get_slides(short=False):
    """Return ordered list of slide specs.  Each is a dict consumed by builders."""

    slides = []

    # ── Phase 1: The Vision ───────────────────────────────────────────────
    SEC1 = "THE VISION"

    slides.append({"type": "title"})

    slides.append({
        "type": "standard",
        "title": "What is the Eat Jersey Challenge?",
        "section": SEC1,
        "bullets": [
            "Visit a restaurant in every one of New Jersey's 564 municipalities",
            "Share each visit on Instagram, TikTok, and Facebook",
            "Document the journey — every borough, city, township, town, and village",
            "564 towns across 21 counties — a massive undertaking",
            "Started with a spreadsheet… but we needed something better",
        ],
        "notes": (
            "Set the stage. Explain the challenge: we want to eat at a restaurant "
            "in every single municipality in NJ. 564 is a LOT. We needed a way to "
            "track progress visually, not just rows in a spreadsheet."
        ),
    })

    slides.append({
        "type": "standard",
        "title": "The Vision: An Interactive Map",
        "section": SEC1,
        "bullets": [
            "Colour-coded map: gold = visited, grey = unvisited, orange = queued",
            "Click any town to see restaurant, date, photos, and social links",
            "Show progress at a glance — how much of NJ have we covered?",
            "Share it publicly so friends and followers can explore",
            "Replace the spreadsheet with something people actually want to look at",
        ],
        "notes": (
            "The vision was simple: instead of a spreadsheet, build an interactive "
            "map where you can see every town coloured by status and click for details. "
            "Something you can share with people and they'll actually explore."
        ),
        "placeholder_label": "Screenshot: Live map with visited/unvisited towns",
    })

    if not short:
        slides.append({
            "type": "demo",
            "demo_title": "The Live Map",
            "section": SEC1,
            "instructions": [
                "Show the map at full zoom — gold vs grey towns",
                "Click a visited town to show the popup with restaurant, date, social links",
                "Toggle the detail view to show queued (orange) and pre-challenge (purple)",
                "Show the legend toggle",
            ],
            "notes": (
                "DEMO: Open the live map. Show the colour coding. Click on a visited "
                "town and walk through the popup: restaurant name, date, thumbnail, "
                "social media icons. Toggle detail mode to show queued and pre-challenge."
            ),
        })

    # ── Phase 2: Technology Choices ────────────────────────────────────────
    SEC2 = "TECHNOLOGY CHOICES"

    if short:
        slides.append({
            "type": "standard",
            "title": "Technology Stack",
            "section": SEC2,
            "bullets": [
                "Frontend: Vanilla JavaScript — zero frameworks, zero build tools",
                "Mapping: Leaflet.js + OpenStreetMap tiles — free and open source",
                "Data: Static JSON file committed to Git — no database needed",
                "Hosting: GitHub Pages — free, automatic deployment on push",
                "Scripts: Node.js + Python for data import and automation",
                "CI/CD: GitHub Actions for automated thumbnail fetching",
                "Why static? Low maintenance, version-controlled data, no backend costs",
            ],
            "notes": (
                "Walk through the tech stack. Emphasize the simplicity: no React, "
                "no database, no server. Everything is a static file served from GitHub. "
                "Data changes are Git commits. This is intentional — low maintenance."
            ),
        })
    else:
        slides.append({
            "type": "standard",
            "title": "Technology Stack Overview",
            "section": SEC2,
            "two_col": {
                "left": [
                    "Frontend: Vanilla JavaScript (ES5+)",
                    "Mapping: Leaflet.js 1.9 + OpenStreetMap",
                    "Styling: CSS3 — no preprocessors",
                    "Icons: Font Awesome 6.7 (social brands)",
                ],
                "right": [
                    "Data: Static JSON committed to Git",
                    "Hosting: GitHub Pages",
                    "Scripts: Node.js + Python 3",
                    "CI/CD: GitHub Actions",
                ],
            },
            "notes": (
                "Quick overview of the full stack. Key takeaway: this is a zero-dependency "
                "frontend with no build step. Data lives in Git. Hosting is free."
            ),
        })

        slides.append({
            "type": "standard",
            "title": "Why No Framework?",
            "section": SEC2,
            "bullets": [
                "Zero frontend dependencies — no npm install for the website itself",
                "No build step — edit HTML/JS/CSS and push; GitHub Pages does the rest",
                "Static-first: all data is a JSON file versioned in Git",
                "Lower maintenance burden — no dependency updates, no breaking changes",
                "Perfect for a personal project that needs to last for years",
                "When you don't need React, don't use React",
            ],
            "notes": (
                "Justify the choice. This isn't an anti-framework talk — it's about picking "
                "the right tool. A personal mapping project doesn't need a SPA framework. "
                "Vanilla JS + static hosting = years of zero-maintenance operation."
            ),
        })

    slides.append({
        "type": "standard",
        "title": "How to Make a Map",
        "section": SEC2,
        "bullets": [
            "Leaflet.js: lightweight open-source mapping library",
            "OpenStreetMap tiles: free map imagery, no API key required",
            "GeoJSON: NJ municipality boundary polygons from US Census Bureau",
            "Each polygon keyed by GEOID — a unique 10-digit Census identifier",
            "Colour each polygon based on visit status from our JSON data",
            "Add click handlers to open rich popups with town details",
        ],
        "notes": (
            "Explain the three ingredients: a mapping library (Leaflet), map tiles "
            "(OpenStreetMap — free), and boundary polygons (GeoJSON from Census). "
            "The GEOID is the magic key that joins geographic shapes to our data."
        ),
        "placeholder_label": "Diagram: Leaflet + OSM Tiles + GeoJSON = Interactive Map",
    })

    slides.append({
        "type": "standard",
        "title": "The Data Model",
        "section": SEC2,
        "bullets": [
            "municipalities.json: flat object keyed by 10-digit GEOID",
            "Each entry: name, county, townType, status, visitNumber",
            "Restaurant info: restaurantName, dateVisited",
            "Links array: [{category, platform, label, url}] — flexible & extensible",
            "Status values: visited | unvisited | queued | pre-challenge",
            "Schema evolved from simple boolean 'visited' to rich multi-status model",
        ],
        "notes": (
            "Walk through the data model. Show a sample JSON entry. Emphasize: "
            "GEOID-based keys eliminate name ambiguity (there are multiple 'Washington' "
            "towns). The links array is flexible — any platform, any category."
        ),
    })

    slides.append({
        "type": "standard",
        "title": "From Firebase to Static JSON",
        "section": SEC2,
        "bullets": [
            "Originally started with Firebase Firestore as the database",
            "Realized: the data only changes when WE update it — no real-time needed",
            "Removed Firebase entirely (visible in early commit history)",
            "Static JSON = zero cost, zero auth, zero backend maintenance",
            "Git versioning gives us full history of every data change for free",
            "Lesson: don't add infrastructure you don't need",
        ],
        "notes": (
            "This is an important architecture lesson. We started with Firebase "
            "because it seemed 'proper' — but we only update data weekly. A static "
            "JSON file committed to Git is simpler, cheaper, and gives version history."
        ),
    })

    if not short:
        slides.append({
            "type": "demo",
            "demo_title": "Admin Panel & Data Flow",
            "section": SEC2,
            "instructions": [
                "Open admin.html — show the searchable town list",
                "Click a town — show the edit modal with status, restaurant, links",
                "Edit a field, save, then export JSON",
                "Show the downloaded municipalities.json",
            ],
            "notes": (
                "DEMO: Walk through the admin panel. Search for a town, open it, "
                "show the edit form. Explain: changes are in-memory only until you "
                "Export JSON. This exported file gets committed to Git."
            ),
        })

    # ── Phase 3: Data Workflows ────────────────────────────────────────────
    SEC3 = "DATA WORKFLOWS"

    if short:
        slides.append({
            "type": "standard",
            "title": "Importing Data from Excel",
            "section": SEC3,
            "bullets": [
                "Generate a colour-coded Excel template from municipalities.json",
                "Edit in Excel — bulk updates to statuses, restaurants, dates, links",
                "Export as CSV, then run: npm run import path/to/file.csv",
                "Script merges CSV data back into municipalities.json",
                "Preserves 'orphaned' links added via admin panel (not in Excel columns)",
                "Handles blank cells carefully — doesn't accidentally erase data",
            ],
            "notes": (
                "Explain the Excel workflow. The Python script generates a rich Excel "
                "template with formulas, colour coding, and summary sheets. Edit in "
                "Excel for bulk operations, then import back. The CSV parser is "
                "robust — handles quoted fields, preserves orphaned data."
            ),
        })
    else:
        slides.append({
            "type": "standard",
            "title": "Importing Data from Excel",
            "section": SEC3,
            "bullets": [
                "Existing data lived in an Excel spreadsheet before the website",
                "Built a CSV-to-JSON import pipeline: npm run import path/to/file.csv",
                "Custom CSV parser handles quoted fields with embedded commas",
                "Preserves 'orphaned' links — data added via admin that isn't in Excel",
                "Blank cell handling: won't accidentally erase existing data on import",
            ],
            "notes": (
                "The challenge: we had data in Excel already. We needed to get it into "
                "our JSON format. The import script is smart — it merges, not replaces. "
                "Orphaned links (added via admin panel) are preserved and flagged with warnings."
            ),
        })

        slides.append({
            "type": "standard",
            "title": "Generating an Excel Template",
            "section": SEC3,
            "bullets": [
                "Python script: npm run gen-template",
                "Generates a full Excel workbook from municipalities.json",
                "Pre-populated with all 564 towns and existing data",
                "Colour-coded columns: identity, visit info, restaurant, social platforms",
                "Summary sheet with live formulas showing progress by county",
                "Auto-scales platform columns (2–5 slots per platform based on usage)",
            ],
            "notes": (
                "Show the generated Excel file. Highlight: it's not just a dump — it has "
                "colour-coded columns, a summary sheet with formulas, and auto-scaled "
                "platform columns. The template is the bridge between Excel users and the website."
            ),
        })

    if short:
        slides.append({
            "type": "standard",
            "title": "Social Media & Thumbnail Automation",
            "section": SEC3,
            "bullets": [
                "Instagram: export posts via PostFox → Python script matches to towns",
                "Facebook: Graph API fetch → Python script matches to towns",
                "Both: normalize labels, preserve existing data, create backups",
                "Thumbnails: Instagram shortcodes → download → resize to 320px WebP",
                "GitHub Actions: push PostFox export → auto-fetch new thumbnails",
                "Fully automated CI/CD for a static site!",
            ],
            "notes": (
                "Two bulk import scripts pull social media posts and match them to towns "
                "by visit number. Thumbnails are fetched from Instagram and auto-committed "
                "by GitHub Actions when a new PostFox export is pushed."
            ),
        })
    else:
        slides.append({
            "type": "standard",
            "title": "Social Media Bulk Import",
            "section": SEC3,
            "bullets": [
                "Instagram: PostFox browser extension exports all posts as JSON",
                "Python script parses export, matches posts to towns by visit number",
                "Facebook: Graph API fetches page posts automatically",
                "Both scripts: normalize labels ('NJ Town #N: Name'), create backups",
                "Output populates Excel columns → re-import via CSV pipeline",
                "Result: social links for every visited town, added in bulk",
            ],
            "notes": (
                "Walk through the Instagram workflow: PostFox extension → JSON export → "
                "Python script matches by visit number → populates Excel → CSV import. "
                "Same pattern for Facebook but using the Graph API directly."
            ),
        })

        slides.append({
            "type": "standard",
            "title": "Thumbnail Automation with GitHub Actions",
            "section": SEC3,
            "bullets": [
                "Thumbnails: Instagram post images shown in map popups",
                "Script extracts shortcodes from Instagram URLs → downloads images",
                "Images resized to 320px WebP using Sharp (Node.js)",
                "GitHub Actions workflow: triggered when PostFox export is committed",
                "Automatically fetches missing thumbnails and commits new .webp files",
                "Idempotent: safe to re-run — already-downloaded images are skipped",
            ],
            "notes": (
                "Show the GitHub Actions workflow file. Explain: when someone pushes a "
                "new PostFox export, the CI automatically downloads any missing thumbnail "
                "images and commits them. Zero manual work for image management."
            ),
            "placeholder_label": "Screenshot: GitHub Actions workflow run",
        })

    slides.append({
        "type": "standard",
        "title": "Edge Cases & Data Quirks",
        "section": SEC3,
        "bullets": [
            "Ambiguous town names: 4 towns named 'Washington' across different counties",
            "Solution: maintain an ambiguousNames set, append county to Wikipedia links",
            "Schema migration: v1 had boolean 'visited' → v2 has 4-status model",
            "Migration script: scripts/migrate-schema-v2.js for backward compatibility",
            "SRI hash mismatch: Leaflet CDN integrity hash changed unexpectedly",
            "GeoJSON size: full NJ polygon file is large — lazy loading considerations",
        ],
        "notes": (
            "Interesting edge cases we hit. The ambiguous names problem was subtle — "
            "multiple towns share the same name in different counties. Schema migration "
            "shows how data models evolve. The SRI hash issue was a fun debugging exercise."
        ),
    })

    if not short:
        slides.append({
            "type": "demo",
            "demo_title": "Excel Import Workflow",
            "section": SEC3,
            "instructions": [
                "Show the generated Excel template (colour-coded columns, summary sheet)",
                "Make a change in Excel — mark a town as visited",
                "Export as CSV",
                "Run: npm run import path/to/file.csv",
                "Show the updated municipalities.json diff",
            ],
            "notes": (
                "DEMO: Open the Excel template. Point out colour coding and formulas. "
                "Make a change, export CSV, run the import. Show the JSON diff to prove "
                "it worked. Highlight orphaned link preservation."
            ),
        })

    # ── Phase 4: Building with GitHub Copilot ──────────────────────────────
    SEC4 = "BUILDING WITH GITHUB COPILOT"

    slides.append({
        "type": "standard",
        "title": "Vibe Coding: The Starting Point",
        "section": SEC4,
        "bullets": [
            "The very first commit: 'initial set of vibe coded files'",
            "GitHub Copilot (Agent mode) generated the entire first version",
            "HTML, CSS, JavaScript, even the initial data structure",
            "From idea to working prototype in a single session",
            "Not perfect — but a working starting point to iterate on",
            "This is the promise of AI-assisted development",
        ],
        "notes": (
            "The origin story. The second commit in the repo is literally called "
            "'initial set of vibe coded files'. Copilot's Agent mode generated the "
            "entire first version of the map. It worked! Not perfectly, but it was "
            "a functional prototype that we could iterate on."
        ),
    })

    if short:
        slides.append({
            "type": "standard",
            "title": "Plan vs Agent vs Ask Mode",
            "section": SEC4,
            "bullets": [
                "Plan Mode: researches first, proposes changes across files, asks for approval",
                "  → Best for: architecture decisions, multi-file refactors, new features",
                "Agent Mode: implements directly — edits files, runs commands, iterates",
                "  → Best for: writing code, running scripts, debugging, building features",
                "Ask Mode: answers questions, explains code — read-only",
                "  → Best for: understanding unfamiliar code, quick questions, learning",
                "The handoff: Plan the approach → Agent builds it → Ask to debug",
            ],
            "notes": (
                "Key lesson: knowing which mode to use when. Plan mode is great for "
                "big-picture decisions. Agent mode is the workhorse. Ask mode is your "
                "rubber duck. The best workflow chains them together."
            ),
        })
    else:
        slides.append({
            "type": "standard",
            "title": "Three Modes of Copilot",
            "section": SEC4,
            "bullets": [
                "Plan Mode: research → propose → approve → execute",
                "Agent Mode: autonomous implementation with tool access",
                "Ask Mode: conversational Q&A, read-only",
                "Each mode has a sweet spot — choosing wisely saves time",
                "The real power: chaining modes in a workflow",
            ],
            "notes": (
                "Introduce the three modes. Set up the next three slides that go deeper "
                "into each one with concrete examples from this project."
            ),
        })

        slides.append({
            "type": "standard",
            "title": "Plan Mode Deep Dive",
            "section": SEC4,
            "bullets": [
                "Best for: architecture decisions and multi-file changes",
                "Researches the codebase before proposing anything",
                "Shows a step-by-step plan you can review and approve",
                "Example: planning the schema migration from v1 to v2",
                "Example: designing the Excel template generation script",
                "Tip: use Plan mode when you're not sure HOW to approach a problem",
            ],
            "notes": (
                "Plan mode excels when you need to think before you act. It reads "
                "multiple files, understands relationships, and proposes a coherent plan. "
                "We used it for the schema migration — it identified all the files that "
                "needed to change and proposed the migration script."
            ),
        })

        slides.append({
            "type": "standard",
            "title": "Agent Mode Deep Dive",
            "section": SEC4,
            "bullets": [
                "Best for: writing code, running scripts, iterating",
                "Can edit files, run terminal commands, and respond to errors",
                "Generated the bulk import scripts (Python and Node.js)",
                "Built the admin panel UI from a description",
                "Iterative: 'This doesn't work' → it reads the error → it fixes it",
                "Tip: give it clear, specific instructions — don't be vague",
            ],
            "notes": (
                "Agent mode is the workhorse. It generated most of the scripts in this "
                "project. The key insight: it can run code and see errors, so it iterates "
                "like a developer would. Give it specific instructions and let it work."
            ),
        })

        slides.append({
            "type": "standard",
            "title": "Ask Mode Deep Dive",
            "section": SEC4,
            "bullets": [
                "Best for: understanding code, quick questions, debugging",
                "Read-only — doesn't change anything",
                "'What does this function do?' → instant explanation",
                "'Why is this polygon not rendering?' → targeted debugging help",
                "Great for onboarding: 'Explain this codebase to me'",
                "Tip: use Ask to understand before using Agent to change",
            ],
            "notes": (
                "Ask mode is your rubber duck. Use it to understand unfamiliar code "
                "or debug issues. It's read-only so it's safe — no accidental changes. "
                "Great pattern: Ask to understand → Plan to design → Agent to build."
            ),
        })

    slides.append({
        "type": "standard",
        "title": "Copilot Pain Points & Issues",
        "section": SEC4,
        "bullets": [
            "VS Code crashes: lost context window and conversation history",
            "No way to recover a conversation after a crash",
            "Context window limits: long sessions degrade response quality",
            "Copilot sometimes 'forgets' earlier decisions in long conversations",
            "Generated code needs review — it can introduce subtle bugs",
            "Workaround: break work into smaller, focused sessions",
        ],
        "notes": (
            "Be honest about the challenges. VS Code crashes were the most painful — "
            "losing an entire conversation with context and decisions. The context window "
            "limit is real: after ~30 minutes of back-and-forth, quality drops. "
            "Our workaround: shorter sessions, commit often, use memory files."
        ),
    })

    slides.append({
        "type": "standard",
        "title": "Tips & Best Practices",
        "section": SEC4,
        "bullets": [
            "Use memory files: .github/copilot-instructions.md for project context",
            "Break large tasks into small, focused sessions",
            "Commit frequently — you can always revert",
            "Be specific with prompts: 'Add a gold border to visited polygons' > 'Make it look better'",
            "Review every line of generated code — trust but verify",
            "Plan mode first, then Agent — don't jump straight to implementation",
        ],
        "notes": (
            "Practical tips from real experience. The memory files tip is huge — "
            "it gives Copilot persistent context about your project. Being specific "
            "with prompts is the single biggest quality lever. And always review the code."
        ),
    })

    if not short:
        slides.append({
            "type": "demo",
            "demo_title": "GitHub Copilot in Action",
            "section": SEC4,
            "instructions": [
                "Show a real Copilot conversation in VS Code",
                "Demonstrate Plan mode: ask it to design a new feature",
                "Switch to Agent mode: have it implement the plan",
                "Show how it reads errors and self-corrects",
                "Show a memory file and how it provides context",
            ],
            "notes": (
                "DEMO: Live Copilot usage. Start with Plan mode — ask it to design "
                "something. Review the plan. Switch to Agent and have it build. "
                "Show the iterative error-correction cycle. Show a memory file."
            ),
        })

    # ── Phase 5: The Journey & Wrap-up ─────────────────────────────────────
    SEC5 = "THE JOURNEY"

    slides.append({
        "type": "standard",
        "title": "Evolution Timeline",
        "section": SEC5,
        "bullets": [
            "Initial commit → 'vibe coded files' — Copilot generates v1",
            "Remove Firebase → pivot to static JSON architecture",
            "Schema migration → boolean visited → 4-status model",
            "Admin panel → visual editing without touching JSON",
            "Excel import pipeline → bulk data management",
            "Social media imports → Instagram & Facebook automation",
            "Thumbnail automation → GitHub Actions CI/CD",
            "Iterative popup refinement → 10+ styling commits",
        ],
        "notes": (
            "Walk through the project timeline using commit history milestones. "
            "Each step built on the previous one. Emphasize: this wasn't planned "
            "up front — it evolved organically based on needs."
        ),
    })

    if not short:
        slides.append({
            "type": "standard",
            "title": "The Power of Iterative Refinement",
            "section": SEC5,
            "bullets": [
                "10+ commits just tweaking popup styles and thumbnail display",
                "Font sizes, padding, overlay effects, hover states — small details",
                "Each commit is small, focused, and reversible",
                "Don't try to get it perfect the first time — iterate",
                "Git makes this safe: every change is a commit you can revert",
                "AI-assisted iteration: 'Make the font smaller' → instant result",
            ],
            "notes": (
                "The commit history tells a story of relentless refinement. Over 10 "
                "commits just for popup styling. This is normal and healthy — small "
                "incremental improvements compound. Git makes it safe to experiment."
            ),
        })

    slides.append({
        "type": "standard",
        "title": "What I'd Do Differently",
        "section": SEC5,
        "bullets": [
            "Start with the data model — get that right first",
            "Set up the Excel pipeline earlier — bulk entry saves hours",
            "Use TypeScript or JSDoc for better Copilot suggestions",
            "Add basic tests — even for a personal project",
            "Document decisions as you go — not after the fact",
            "Shorter Copilot sessions from day one to avoid context loss",
        ],
        "notes": (
            "Honest retrospective. The data model evolved twice — getting it right "
            "first would have saved a migration. The Excel pipeline was a late addition "
            "that would have been valuable from the start. TypeScript would give "
            "Copilot better type context for suggestions."
        ),
    })

    slides.append({
        "type": "standard",
        "title": "By the Numbers",
        "section": SEC5,
        "bullets": [
            "564 municipalities in New Jersey",
            "[X] towns visited so far",
            "[X] commits in the repository",
            "0 frameworks, 0 build tools, 0 monthly hosting cost",
            "~2,500 lines of JavaScript (map.js + admin.js + config.js)",
            "6 automation scripts (Node.js + Python)",
            "1 GitHub Actions workflow",
        ],
        "notes": (
            "Fun slide with project stats. Update the [X] placeholders with real "
            "numbers before presenting. The zero-cost angle is compelling — "
            "this entire project runs for free on GitHub Pages."
        ),
    })

    if not short:
        slides.append({
            "type": "standard",
            "title": "What's Next?",
            "section": SEC5,
            "bullets": [
                "Continue visiting — 564 towns is a multi-year journey",
                "Add county-level progress visualization",
                "Implement search/filter on the public map",
                "Mobile-friendly improvements",
                "User submissions? Community contributions?",
                "The website grows with the challenge",
            ],
            "notes": (
                "Share future plans. The project isn't done — it evolves as the "
                "challenge continues. Open it up for ideas from the audience."
            ),
        })

    slides.append({"type": "thankyou"})

    return slides


# ── Main build ─────────────────────────────────────────────────────────────

def build_presentation(short=False):
    slide_specs = get_slides(short=short)
    total = len(slide_specs)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for idx, spec in enumerate(slide_specs, start=1):
        t = spec["type"]
        if t == "title":
            make_title_slide(prs, idx, total)
        elif t == "thankyou":
            make_thank_you_slide(prs, idx, total)
        elif t == "demo":
            make_demo_slide(
                prs,
                demo_title=spec["demo_title"],
                instructions=spec["instructions"],
                notes=spec["notes"],
                section=spec["section"],
                slide_num=idx,
                total=total,
            )
        elif t == "standard":
            make_standard_slide(
                prs,
                title=spec["title"],
                bullets=spec.get("bullets"),
                notes=spec.get("notes", ""),
                section=spec.get("section"),
                slide_num=idx,
                total=total,
                placeholder_label=spec.get("placeholder_label"),
                two_col=spec.get("two_col"),
            )

    suffix = "-Short" if short else ""
    filename = f"EatJerseyChallenge-Presentation{suffix}.pptx"
    out_path = os.path.join(os.path.dirname(__file__), "..", filename)
    prs.save(out_path)
    abs_path = os.path.abspath(out_path)
    print(f"✅  Saved {filename}  ({total} slides)")
    print(f"    {abs_path}")
    return abs_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate the Eat Jersey Challenge PowerPoint presentation.")
    parser.add_argument("--short", action="store_true",
                        help="Generate the 15-min condensed version (~18 slides)")
    parser.add_argument("--both", action="store_true",
                        help="Generate both the full and short versions")
    args = parser.parse_args()

    if args.both:
        build_presentation(short=False)
        build_presentation(short=True)
    else:
        build_presentation(short=args.short)


if __name__ == "__main__":
    main()
